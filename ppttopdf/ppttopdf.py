import os
import comtypes.client
import sys

def ppt_to_pdf(input_file_path, output_file_path):
    """
    将 PowerPoint 文件转换为 PDF。

    :param input_file_path: 输入的 PowerPoint 文件 (.ppt 或 .pptx) 的绝对路径。
    :param output_file_path: 输出的 PDF 文件的绝对路径。
    """
    # 获取文件的绝对路径，这对于 comtypes 至关重要
    input_file_path = os.path.abspath(input_file_path)
    output_file_path = os.path.abspath(output_file_path)

    powerpoint = None
    try:
        # 启动 PowerPoint 应用程序
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        # 设置为在后台运行，不显示窗口
        powerpoint.Visible = 1

        # 打开输入的 PowerPoint 文件
        presentation = powerpoint.Presentations.Open(input_file_path)

        # 将演示文稿另存为 PDF
        # formatType=32 表示保存为 PDF 格式
        presentation.SaveAs(output_file_path, 32)
        print(f"成功将 '{input_file_path}' 转换为 '{output_file_path}'")

        # 关闭演示文稿
        presentation.Close()

    except Exception as e:
        print(f"转换过程中发生错误: {e}")
    finally:
        # 确保 PowerPoint 应用程序被关闭
        if powerpoint:
            powerpoint.Quit()

# --- 使用示例 ---
if __name__ == '__main__':
    # 检查是否在 Windows 上运行
    if sys.platform != 'win32':
        print("错误：此脚本需要 Windows 操作系统和 Microsoft PowerPoint。")
    else:
        # 设置输入和输出文件路径
        # 请将 "Your_Presentation.pptx" 替换为您的实际文件名
        input_ppt = "第1章_绪论.ppt"
        output_pdf = "Converted_Presentation.pdf"

        # 检查输入文件是否存在
        if not os.path.exists(input_ppt):
            print(f"错误: 输入文件 '{input_ppt}' 未找到。")
            # 为了演示，创建一个虚拟的PPTX文件
            print("正在创建一个临时的PPTX文件用于演示...")
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = "Hello, World!"
            prs.save(input_ppt)
            print(f"已创建 '{input_ppt}'。")

        # 调用转换函数
        ppt_to_pdf(input_ppt, output_pdf)